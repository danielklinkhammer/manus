import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Pfad zur Vorlage
template_path = '/home/ubuntu/bootstrap_presentation/Vorlesung 1 KOPR Kopie.pptx'
output_path = '/home/ubuntu/bootstrap_presentation/Bootstrap_Kurs_Präsentation.pptx'

# Präsentation öffnen
prs = Presentation(template_path)

# Funktion zum Hinzufügen einer Folie mit Titel und Inhalt
def add_title_content_slide(title, content, layout_index=1):
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titel hinzufügen
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Inhalt hinzufügen
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.text = content
    
    return slide

# Funktion zum Hinzufügen einer Folie mit Titel und Aufzählungspunkten
def add_title_bullets_slide(title, bullets, layout_index=1):
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titel hinzufügen
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Aufzählungspunkte hinzufügen
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
    
    return slide

# Funktion zum Hinzufügen einer Titelfolie
def add_title_slide(title, subtitle="", layout_index=0):
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    
    # Titel hinzufügen
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Untertitel hinzufügen, wenn vorhanden
    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
    
    return slide

# Titelfolie hinzufügen
add_title_slide("Bootstrap für UX Design", "Ein umfassender Kurs für Studierende")

# Agenda hinzufügen
add_title_bullets_slide("Agenda", [
    "Einführung in Bootstrap",
    "HTML/CSS Grundlagen (Wiederholung)",
    "Bootstrap Grid-System",
    "Bootstrap-Komponenten",
    "Responsives Design mit Bootstrap",
    "Figma Dev Mode zu Bootstrap",
    "Praktische Übungen",
    "Ressourcen und weiterführende Materialien"
])

# Einführung in Bootstrap
add_title_bullets_slide("Was ist Bootstrap?", [
    "Frontend-Framework für die Webentwicklung",
    "Entwickelt von Twitter, 2011 veröffentlicht",
    "Aktuell: Bootstrap 5 (ohne jQuery-Abhängigkeit)",
    "Mobile-first Ansatz",
    "Umfangreiche Komponenten-Bibliothek",
    "Responsives Grid-System",
    "Konsistentes Design auf allen Geräten"
])

add_title_bullets_slide("Vorteile von Bootstrap", [
    "Schnelle Entwicklung von responsiven Websites",
    "Konsistentes Design und Benutzerfreundlichkeit",
    "Große Community und umfangreiche Dokumentation",
    "Kompatibilität mit allen modernen Browsern",
    "Anpassbar durch Sass-Variablen",
    "Regelmäßige Updates und Verbesserungen",
    "Ideal für Prototyping und Produktiventwicklung"
])

add_title_bullets_slide("Bootstrap vs. andere Frameworks", [
    "Bootstrap: Umfassend, gut dokumentiert, große Community",
    "Tailwind CSS: Utility-first, hochgradig anpassbar",
    "Foundation: Ähnlich wie Bootstrap, mehr Enterprise-fokussiert",
    "Bulma: Leichtgewichtig, modern, ohne JavaScript",
    "Material UI: Implementiert Google's Material Design",
    "Bootstrap ist ideal für Einsteiger und schnelle Entwicklung"
])

# HTML/CSS Grundlagen
add_title_slide("HTML/CSS Grundlagen", "Wiederholung der wichtigsten Konzepte")

add_title_bullets_slide("HTML Grundlagen", [
    "HTML = HyperText Markup Language",
    "Struktur einer Webseite definieren",
    "Elemente werden durch Tags dargestellt: <tag>Inhalt</tag>",
    "Wichtige Tags: html, head, body, div, p, h1-h6, a, img, ul/ol, li",
    "Semantische Tags: header, nav, main, section, article, footer",
    "Attribute erweitern Tags: class, id, src, href, alt"
])

add_title_bullets_slide("CSS Grundlagen", [
    "CSS = Cascading Style Sheets",
    "Styling und Layout einer Webseite",
    "Selektoren: Element, Klasse (.klasse), ID (#id)",
    "Eigenschaften: color, font-size, margin, padding, display",
    "Box-Modell: content, padding, border, margin",
    "Flexbox und Grid für Layout",
    "Media Queries für Responsivität"
])

add_title_bullets_slide("Wichtige CSS-Konzepte für Bootstrap", [
    "Flexbox: Flexible Box Layout für eindimensionale Layouts",
    "CSS Grid: Zweidimensionales Layoutsystem",
    "Media Queries: Anpassung des Designs an verschiedene Bildschirmgrößen",
    "Box-Modell: Verständnis von margin, padding, border",
    "Positionierung: relative, absolute, fixed, sticky",
    "Bootstrap baut auf diesen Konzepten auf und vereinfacht sie"
])

# Bootstrap Grid-System
add_title_slide("Bootstrap Grid-System", "Das Fundament für responsive Layouts")

add_title_bullets_slide("Bootstrap Grid-System Grundlagen", [
    "12-Spalten-System für flexible Layouts",
    "Container → Rows → Columns",
    "Responsive Breakpoints: xs, sm, md, lg, xl, xxl",
    "Mobile-first Ansatz",
    "Klassen: .container, .row, .col-*",
    "Automatische Anpassung an Bildschirmgrößen"
])

add_title_bullets_slide("Container-Typen", [
    ".container: Feste Breite, zentriert, responsive",
    ".container-fluid: Volle Breite auf allen Geräten",
    ".container-{breakpoint}: 100% Breite bis zum Breakpoint",
    "Breakpoints: sm (576px), md (768px), lg (992px), xl (1200px), xxl (1400px)"
])

add_title_bullets_slide("Spalten und Zeilen", [
    ".row: Erzeugt eine horizontale Gruppe von Spalten",
    ".col-*: Definiert die Spaltenbreite (1-12)",
    ".col-{breakpoint}-*: Responsive Spaltenbreite",
    "Beispiel: .col-md-6 (50% Breite ab mittleren Geräten)",
    "Auto-Layout: .col oder .col-auto",
    "Verschachtelte Grids möglich (Rows innerhalb von Columns)"
])

add_title_bullets_slide("Responsive Utilities", [
    "Offset: .offset-*-* (Verschieben von Spalten)",
    "Order: .order-*-* (Ändern der Reihenfolge)",
    "Alignment: .align-items-*, .align-self-*",
    "Justify Content: .justify-content-*",
    "Display Utilities: .d-none, .d-{breakpoint}-block",
    "Margin und Padding: .m-*, .p-*, .mx-*, .py-*"
])

# Bootstrap-Komponenten
add_title_slide("Bootstrap-Komponenten", "Fertige UI-Elemente für moderne Websites")

add_title_bullets_slide("Navigationskomponenten", [
    "Navbar: Responsive Navigationsleiste",
    "Nav: Tabs, Pills, vertikale Navigation",
    "Breadcrumb: Navigationspfad",
    "Pagination: Seitennummerierung",
    "Dropdown: Ausklappbare Menüs",
    "Offcanvas: Seitliche Navigationsleiste"
])

add_title_bullets_slide("Inhaltskomponenten", [
    "Cards: Flexible Container für Inhalte",
    "Carousel: Bildslider",
    "Modal: Dialogfenster",
    "Accordion: Ausklappbare Panels",
    "Tabs: Tabbed Interface",
    "List Group: Formatierte Listen"
])

add_title_bullets_slide("Formularkomponenten", [
    "Form Controls: Input, Textarea, Select",
    "Checks & Radios: Checkboxen und Radiobuttons",
    "Input Group: Kombinierte Eingabefelder",
    "Floating Labels: Schwebende Beschriftungen",
    "Form Validation: Validierung von Eingaben",
    "Range: Schieberegler"
])

add_title_bullets_slide("Feedback-Komponenten", [
    "Alert: Benachrichtigungen und Warnungen",
    "Badge: Kleine Zähler oder Labels",
    "Progress: Fortschrittsbalken",
    "Spinner: Ladeanzeigen",
    "Toast: Temporäre Benachrichtigungen",
    "Tooltip & Popover: Informations-Overlays"
])

# Responsives Design mit Bootstrap
add_title_slide("Responsives Design mit Bootstrap", "Optimale Darstellung auf allen Geräten")

add_title_bullets_slide("Responsive Design Prinzipien", [
    "Mobile-first Ansatz: Design für kleine Bildschirme zuerst",
    "Fluid Grids: Relative statt absolute Größen",
    "Flexible Bilder: max-width: 100%",
    "Media Queries: Anpassung an verschiedene Bildschirmgrößen",
    "Viewport Meta-Tag: <meta name=\"viewport\" content=\"width=device-width, initial-scale=1\">",
    "Bootstrap implementiert all diese Prinzipien automatisch"
])

add_title_bullets_slide("Responsive Breakpoints in Bootstrap", [
    "Extra Small (xs): < 576px (Smartphones)",
    "Small (sm): ≥ 576px (Große Smartphones, kleine Tablets)",
    "Medium (md): ≥ 768px (Tablets)",
    "Large (lg): ≥ 992px (Desktops)",
    "Extra Large (xl): ≥ 1200px (Große Desktops)",
    "Extra Extra Large (xxl): ≥ 1400px (Sehr große Displays)"
])

add_title_bullets_slide("Responsive Utilities", [
    "Display: .d-none, .d-{breakpoint}-block",
    "Flexbox: .d-flex, .flex-{breakpoint}-row",
    "Text Alignment: .text-{breakpoint}-center",
    "Spacing: .m-{breakpoint}-*, .p-{breakpoint}-*",
    "Visibility: .visible, .invisible",
    "Responsive Images: .img-fluid"
])

add_title_bullets_slide("Responsive Best Practices", [
    "Testen auf verschiedenen Geräten und Browsern",
    "Verwenden der Browser-Entwicklertools für Simulation",
    "Vermeiden von festen Breiten (px)",
    "Bilder optimieren für schnellere Ladezeiten",
    "Touch-freundliche Elemente (min. 44x44px)",
    "Berücksichtigen von Bildschirmorientierung (Portrait/Landscape)"
])

# Figma Dev Mode zu Bootstrap
add_title_slide("Figma Dev Mode zu Bootstrap", "Vom Design zur Implementierung")

add_title_bullets_slide("Was ist Figma Dev Mode?", [
    "Brücke zwischen Design und Entwicklung",
    "Zugriff auf Design-Tokens und CSS-Eigenschaften",
    "Inspektion von Komponenten und Layouts",
    "Export von Assets und Code",
    "Verbesserte Zusammenarbeit zwischen Designern und Entwicklern",
    "Veröffentlicht 2023 als Nachfolger des Inspect-Modus"
])

add_title_bullets_slide("Figma Dev Mode Features", [
    "Code-Panel: CSS, iOS, Android Code",
    "Variables-Panel: Design-Tokens und Variablen",
    "Layers-Panel: Hierarchie und Struktur",
    "Assets-Panel: Export von Bildern und Icons",
    "Kommentare und Anmerkungen für Entwickler",
    "Integration mit Versionskontrolle"
])

add_title_bullets_slide("Workflow: Figma zu Bootstrap", [
    "1. Design in Figma erstellen",
    "2. Dev Mode aktivieren",
    "3. Design-System auf Bootstrap-Komponenten abbilden",
    "4. CSS-Eigenschaften extrahieren",
    "5. Bootstrap-Klassen identifizieren",
    "6. HTML-Struktur erstellen",
    "7. Bootstrap-Klassen anwenden",
    "8. Anpassungen mit Custom CSS"
])

add_title_bullets_slide("Figma zu Bootstrap Plugins", [
    "Figma to Bootstrap 5 Plugin",
    "Automatische Konvertierung von Figma-Designs zu Bootstrap-Code",
    "Bootstrap 5 UI Kit für Figma",
    "Design mit Bootstrap-Komponenten in Figma",
    "Breakpoints Plugin für responsive Designs",
    "Handoff-Plugins für verbesserte Zusammenarbeit"
])

# Praktische Übungen
add_title_slide("Praktische Übungen", "Hands-on Erfahrung mit Bootstrap")

add_title_bullets_slide("Übung 1: HTML/CSS Grundlagen", [
    "Erstellen einer einfachen Webseite mit HTML und CSS",
    "Semantische HTML-Struktur",
    "CSS-Styling für Farben, Schriftarten, Abstände",
    "Responsive Grundlagen ohne Framework",
    "Ziel: Auffrischung der Grundkenntnisse"
])

add_title_bullets_slide("Übung 2: Bootstrap Grid", [
    "Erstellen eines responsiven Layouts mit Bootstrap Grid",
    "Container, Rows und Columns",
    "Responsive Breakpoints",
    "Verschachtelte Grids",
    "Responsive Utilities (Order, Offset)",
    "Ziel: Verständnis des Grid-Systems"
])

add_title_bullets_slide("Übung 3: Bootstrap Komponenten", [
    "Erstellen einer Produktseite mit Bootstrap-Komponenten",
    "Navbar, Carousel, Cards",
    "Formulare und Buttons",
    "Modals und Tooltips",
    "Accordion und Tabs",
    "Ziel: Anwendung verschiedener Komponenten"
])

add_title_bullets_slide("Übung 4: Responsives Design", [
    "Erstellen einer responsiven Portfolio-Webseite",
    "Mobile-first Ansatz",
    "Anpassung an verschiedene Bildschirmgrößen",
    "Responsive Bilder und Videos",
    "Responsive Utilities",
    "Ziel: Optimierung für alle Geräte"
])

add_title_bullets_slide("Übung 5: Figma zu Bootstrap", [
    "Umsetzung eines Figma-Designs in Bootstrap",
    "Analyse des Designs und Identifikation von Bootstrap-Komponenten",
    "Extraktion von CSS-Eigenschaften aus Figma Dev Mode",
    "Implementierung mit Bootstrap-Klassen",
    "Anpassungen mit Custom CSS",
    "Ziel: Praxisnaher Workflow"
])

add_title_bullets_slide("Abschlussprojekt", [
    "Erstellen einer vollständigen Webseite",
    "Eigenes Design in Figma",
    "Umsetzung mit Bootstrap",
    "Responsive für alle Geräte",
    "Interaktive Elemente",
    "Dokumentation des Entwicklungsprozesses",
    "Präsentation der Ergebnisse"
])

# Ressourcen und weiterführende Materialien
add_title_slide("Ressourcen", "Weiterführende Materialien und Hilfen")

add_title_bullets_slide("Offizielle Dokumentation", [
    "Bootstrap Dokumentation: getbootstrap.com/docs",
    "Figma Help Center: help.figma.com",
    "Visual Studio Code Docs: code.visualstudio.com/docs",
    "MDN Web Docs: developer.mozilla.org"
])

add_title_bullets_slide("Tutorials und Kurse", [
    "Bootstrap 5 Crash Course (YouTube)",
    "Figma für UX/UI Designer (YouTube)",
    "Figma Dev Mode Tutorial (YouTube)",
    "Bootstrap 5: Grundlagen bis Fortgeschrittene (Udemy)",
    "Web Design für UX Designer (Coursera)"
])

add_title_bullets_slide("Tools und Plugins", [
    "Figma to Bootstrap 5 Plugin",
    "Bootstrap 5 UI Kit für Figma",
    "VS Code Erweiterungen: Live Server, Bootstrap Snippets",
    "Bootstrap Themes und Templates",
    "Bootstrap Layout Builder (layoutit.com)"
])

add_title_bullets_slide("Communities und Hilfe", [
    "Stack Overflow: stackoverflow.com/questions/tagged/bootstrap-5",
    "Bootstrap Discord: discord.gg/bZUvakRU2z",
    "Reddit: r/webdev",
    "Figma Community Forum: forum.figma.com",
    "GitHub Issues: github.com/twbs/bootstrap/issues"
])

# Abschlussfolie
add_title_slide("Vielen Dank!", "Fragen und Diskussion")

# Präsentation speichern
prs.save(output_path)

print(f"Präsentation wurde erfolgreich erstellt und unter {output_path} gespeichert.")
