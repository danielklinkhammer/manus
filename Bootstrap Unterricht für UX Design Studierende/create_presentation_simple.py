import os
from pptx import Presentation

# Pfad zur Vorlage und Ausgabedatei
template_path = '/home/ubuntu/bootstrap_presentation/Vorlesung 1 KOPR Kopie.pptx'
output_path = '/home/ubuntu/bootstrap_presentation/Bootstrap_Kurs_Präsentation.pptx'

# Neue Präsentation erstellen (ohne Vorlage zu modifizieren)
prs = Presentation()

# Folien hinzufügen
# Titelfolie
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Bootstrap für UX Design"
subtitle.text = "Ein umfassender Kurs für Studierende"

# Agenda
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Agenda"
content.text = "• Einführung in Bootstrap\n• HTML/CSS Grundlagen (Wiederholung)\n• Bootstrap Grid-System\n• Bootstrap-Komponenten\n• Responsives Design mit Bootstrap\n• Figma Dev Mode zu Bootstrap\n• Praktische Übungen\n• Ressourcen und weiterführende Materialien"

# Einführung in Bootstrap
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Was ist Bootstrap?"
content.text = "• Frontend-Framework für die Webentwicklung\n• Entwickelt von Twitter, 2011 veröffentlicht\n• Aktuell: Bootstrap 5 (ohne jQuery-Abhängigkeit)\n• Mobile-first Ansatz\n• Umfangreiche Komponenten-Bibliothek\n• Responsives Grid-System\n• Konsistentes Design auf allen Geräten"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Vorteile von Bootstrap"
content.text = "• Schnelle Entwicklung von responsiven Websites\n• Konsistentes Design und Benutzerfreundlichkeit\n• Große Community und umfangreiche Dokumentation\n• Kompatibilität mit allen modernen Browsern\n• Anpassbar durch Sass-Variablen\n• Regelmäßige Updates und Verbesserungen\n• Ideal für Prototyping und Produktiventwicklung"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Bootstrap vs. andere Frameworks"
content.text = "• Bootstrap: Umfassend, gut dokumentiert, große Community\n• Tailwind CSS: Utility-first, hochgradig anpassbar\n• Foundation: Ähnlich wie Bootstrap, mehr Enterprise-fokussiert\n• Bulma: Leichtgewichtig, modern, ohne JavaScript\n• Material UI: Implementiert Google's Material Design\n• Bootstrap ist ideal für Einsteiger und schnelle Entwicklung"

# HTML/CSS Grundlagen
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "HTML/CSS Grundlagen"
subtitle.text = "Wiederholung der wichtigsten Konzepte"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "HTML Grundlagen"
content.text = "• HTML = HyperText Markup Language\n• Struktur einer Webseite definieren\n• Elemente werden durch Tags dargestellt: <tag>Inhalt</tag>\n• Wichtige Tags: html, head, body, div, p, h1-h6, a, img, ul/ol, li\n• Semantische Tags: header, nav, main, section, article, footer\n• Attribute erweitern Tags: class, id, src, href, alt"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "CSS Grundlagen"
content.text = "• CSS = Cascading Style Sheets\n• Styling und Layout einer Webseite\n• Selektoren: Element, Klasse (.klasse), ID (#id)\n• Eigenschaften: color, font-size, margin, padding, display\n• Box-Modell: content, padding, border, margin\n• Flexbox und Grid für Layout\n• Media Queries für Responsivität"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Wichtige CSS-Konzepte für Bootstrap"
content.text = "• Flexbox: Flexible Box Layout für eindimensionale Layouts\n• CSS Grid: Zweidimensionales Layoutsystem\n• Media Queries: Anpassung des Designs an verschiedene Bildschirmgrößen\n• Box-Modell: Verständnis von margin, padding, border\n• Positionierung: relative, absolute, fixed, sticky\n• Bootstrap baut auf diesen Konzepten auf und vereinfacht sie"

# Bootstrap Grid-System
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Bootstrap Grid-System"
subtitle.text = "Das Fundament für responsive Layouts"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Bootstrap Grid-System Grundlagen"
content.text = "• 12-Spalten-System für flexible Layouts\n• Container → Rows → Columns\n• Responsive Breakpoints: xs, sm, md, lg, xl, xxl\n• Mobile-first Ansatz\n• Klassen: .container, .row, .col-*\n• Automatische Anpassung an Bildschirmgrößen"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Container-Typen"
content.text = "• .container: Feste Breite, zentriert, responsive\n• .container-fluid: Volle Breite auf allen Geräten\n• .container-{breakpoint}: 100% Breite bis zum Breakpoint\n• Breakpoints: sm (576px), md (768px), lg (992px), xl (1200px), xxl (1400px)"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Spalten und Zeilen"
content.text = "• .row: Erzeugt eine horizontale Gruppe von Spalten\n• .col-*: Definiert die Spaltenbreite (1-12)\n• .col-{breakpoint}-*: Responsive Spaltenbreite\n• Beispiel: .col-md-6 (50% Breite ab mittleren Geräten)\n• Auto-Layout: .col oder .col-auto\n• Verschachtelte Grids möglich (Rows innerhalb von Columns)"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Responsive Utilities"
content.text = "• Offset: .offset-*-* (Verschieben von Spalten)\n• Order: .order-*-* (Ändern der Reihenfolge)\n• Alignment: .align-items-*, .align-self-*\n• Justify Content: .justify-content-*\n• Display Utilities: .d-none, .d-{breakpoint}-block\n• Margin und Padding: .m-*, .p-*, .mx-*, .py-*"

# Bootstrap-Komponenten
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Bootstrap-Komponenten"
subtitle.text = "Fertige UI-Elemente für moderne Websites"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Navigationskomponenten"
content.text = "• Navbar: Responsive Navigationsleiste\n• Nav: Tabs, Pills, vertikale Navigation\n• Breadcrumb: Navigationspfad\n• Pagination: Seitennummerierung\n• Dropdown: Ausklappbare Menüs\n• Offcanvas: Seitliche Navigationsleiste"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Inhaltskomponenten"
content.text = "• Cards: Flexible Container für Inhalte\n• Carousel: Bildslider\n• Modal: Dialogfenster\n• Accordion: Ausklappbare Panels\n• Tabs: Tabbed Interface\n• List Group: Formatierte Listen"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Formularkomponenten"
content.text = "• Form Controls: Input, Textarea, Select\n• Checks & Radios: Checkboxen und Radiobuttons\n• Input Group: Kombinierte Eingabefelder\n• Floating Labels: Schwebende Beschriftungen\n• Form Validation: Validierung von Eingaben\n• Range: Schieberegler"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Feedback-Komponenten"
content.text = "• Alert: Benachrichtigungen und Warnungen\n• Badge: Kleine Zähler oder Labels\n• Progress: Fortschrittsbalken\n• Spinner: Ladeanzeigen\n• Toast: Temporäre Benachrichtigungen\n• Tooltip & Popover: Informations-Overlays"

# Responsives Design mit Bootstrap
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Responsives Design mit Bootstrap"
subtitle.text = "Optimale Darstellung auf allen Geräten"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Responsive Design Prinzipien"
content.text = "• Mobile-first Ansatz: Design für kleine Bildschirme zuerst\n• Fluid Grids: Relative statt absolute Größen\n• Flexible Bilder: max-width: 100%\n• Media Queries: Anpassung an verschiedene Bildschirmgrößen\n• Viewport Meta-Tag: <meta name=\"viewport\" content=\"width=device-width, initial-scale=1\">\n• Bootstrap implementiert all diese Prinzipien automatisch"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Responsive Breakpoints in Bootstrap"
content.text = "• Extra Small (xs): < 576px (Smartphones)\n• Small (sm): ≥ 576px (Große Smartphones, kleine Tablets)\n• Medium (md): ≥ 768px (Tablets)\n• Large (lg): ≥ 992px (Desktops)\n• Extra Large (xl): ≥ 1200px (Große Desktops)\n• Extra Extra Large (xxl): ≥ 1400px (Sehr große Displays)"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Responsive Utilities"
content.text = "• Display: .d-none, .d-{breakpoint}-block\n• Flexbox: .d-flex, .flex-{breakpoint}-row\n• Text Alignment: .text-{breakpoint}-center\n• Spacing: .m-{breakpoint}-*, .p-{breakpoint}-*\n• Visibility: .visible, .invisible\n• Responsive Images: .img-fluid"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Responsive Best Practices"
content.text = "• Testen auf verschiedenen Geräten und Browsern\n• Verwenden der Browser-Entwicklertools für Simulation\n• Vermeiden von festen Breiten (px)\n• Bilder optimieren für schnellere Ladezeiten\n• Touch-freundliche Elemente (min. 44x44px)\n• Berücksichtigen von Bildschirmorientierung (Portrait/Landscape)"

# Figma Dev Mode zu Bootstrap
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Figma Dev Mode zu Bootstrap"
subtitle.text = "Vom Design zur Implementierung"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Was ist Figma Dev Mode?"
content.text = "• Brücke zwischen Design und Entwicklung\n• Zugriff auf Design-Tokens und CSS-Eigenschaften\n• Inspektion von Komponenten und Layouts\n• Export von Assets und Code\n• Verbesserte Zusammenarbeit zwischen Designern und Entwicklern\n• Veröffentlicht 2023 als Nachfolger des Inspect-Modus"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Figma Dev Mode Features"
content.text = "• Code-Panel: CSS, iOS, Android Code\n• Variables-Panel: Design-Tokens und Variablen\n• Layers-Panel: Hierarchie und Struktur\n• Assets-Panel: Export von Bildern und Icons\n• Kommentare und Anmerkungen für Entwickler\n• Integration mit Versionskontrolle"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Workflow: Figma zu Bootstrap"
content.text = "• 1. Design in Figma erstellen\n• 2. Dev Mode aktivieren\n• 3. Design-System auf Bootstrap-Komponenten abbilden\n• 4. CSS-Eigenschaften extrahieren\n• 5. Bootstrap-Klassen identifizieren\n• 6. HTML-Struktur erstellen\n• 7. Bootstrap-Klassen anwenden\n• 8. Anpassungen mit Custom CSS"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Figma zu Bootstrap Plugins"
content.text = "• Figma to Bootstrap 5 Plugin\n• Automatische Konvertierung von Figma-Designs zu Bootstrap-Code\n• Bootstrap 5 UI Kit für Figma\n• Design mit Bootstrap-Komponenten in Figma\n• Breakpoints Plugin für responsive Designs\n• Handoff-Plugins für verbesserte Zusammenarbeit"

# Praktische Übungen
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Praktische Übungen"
subtitle.text = "Hands-on Erfahrung mit Bootstrap"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Übung 1: HTML/CSS Grundlagen"
content.text = "• Erstellen einer einfachen Webseite mit HTML und CSS\n• Semantische HTML-Struktur\n• CSS-Styling für Farben, Schriftarten, Abstände\n• Responsive Grundlagen ohne Framework\n• Ziel: Auffrischung der Grundkenntnisse"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Übung 2: Bootstrap Grid"
content.text = "• Erstellen eines responsiven Layouts mit Bootstrap Grid\n• Container, Rows und Columns\n• Responsive Breakpoints\n• Verschachtelte Grids\n• Responsive Utilities (Order, Offset)\n• Ziel: Verständnis des Grid-Systems"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Übung 3: Bootstrap Komponenten"
content.text = "• Erstellen einer Produktseite mit Bootstrap-Komponenten\n• Navbar, Carousel, Cards\n• Formulare und Buttons\n• Modals und Tooltips\n• Accordion und Tabs\n• Ziel: Anwendung verschiedener Komponenten"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Übung 4: Responsives Design"
content.text = "• Erstellen einer responsiven Portfolio-Webseite\n• Mobile-first Ansatz\n• Anpassung an verschiedene Bildschirmgrößen\n• Responsive Bilder und Videos\n• Responsive Utilities\n• Ziel: Optimierung für alle Geräte"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Übung 5: Figma zu Bootstrap"
content.text = "• Umsetzung eines Figma-Designs in Bootstrap\n• Analyse des Designs und Identifikation von Bootstrap-Komponenten\n• Extraktion von CSS-Eigenschaften aus Figma Dev Mode\n• Implementierung mit Bootstrap-Klassen\n• Anpassungen mit Custom CSS\n• Ziel: Praxisnaher Workflow"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Abschlussprojekt"
content.text = "• Erstellen einer vollständigen Webseite\n• Eigenes Design in Figma\n• Umsetzung mit Bootstrap\n• Responsive für alle Geräte\n• Interaktive Elemente\n• Dokumentation des Entwicklungsprozesses\n• Präsentation der Ergebnisse"

# Ressourcen und weiterführende Materialien
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Ressourcen"
subtitle.text = "Weiterführende Materialien und Hilfen"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Offizielle Dokumentation"
content.text = "• Bootstrap Dokumentation: getbootstrap.com/docs\n• Figma Help Center: help.figma.com\n• Visual Studio Code Docs: code.visualstudio.com/docs\n• MDN Web Docs: developer.mozilla.org"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Tutorials und Kurse"
content.text = "• Bootstrap 5 Crash Course (YouTube)\n• Figma für UX/UI Designer (YouTube)\n• Figma Dev Mode Tutorial (YouTube)\n• Bootstrap 5: Grundlagen bis Fortgeschrittene (Udemy)\n• Web Design für UX Designer (Coursera)"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Tools und Plugins"
content.text = "• Figma to Bootstrap 5 Plugin\n• Bootstrap 5 UI Kit für Figma\n• VS Code Erweiterungen: Live Server, Bootstrap Snippets\n• Bootstrap Themes und Templates\n• Bootstrap Layout Builder (layoutit.com)"

slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Communities und Hilfe"
content.text = "• Stack Overflow: stack<response clipped><NOTE>To save on context only part of this file has been shown to you. You should retry this tool after you have searched inside the file with `grep -n` in order to find the line numbers of what you are looking for.</NOTE>